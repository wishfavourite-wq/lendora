# -*- coding: utf-8 -*-
"""Builds Lendora_Defense_Presentation.pptx from slides_data.SLIDES."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from slides_data import SLIDES

FOREST   = RGBColor(0x0A, 0x3D, 0x2E)
FOREST_D = RGBColor(0x06, 0x28, 0x1E)
COPPER   = RGBColor(0xC8, 0x79, 0x41)
CREAM    = RGBColor(0xFA, 0xF6, 0xF0)
INK      = RGBColor(0x22, 0x26, 0x24)
INK_LT   = RGBColor(0x5A, 0x60, 0x5C)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
PH_BG    = RGBColor(0xEF, 0xEC, 0xE6)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def fill_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        no_line(sp)
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
    return sp


def dashed_rect(slide, l, t, w, h, fill, line):
    sp = rect(slide, l, t, w, h, fill=fill, line=line, line_w=Pt(1.75))
    ln = sp.line._get_or_add_ln()
    dash = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.append(dash)
    return sp


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def set_run(p, text, size, color, bold=False, italic=False, font='Calibri'):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def header_bar(slide, num, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=FOREST)
    rect(slide, 0, Inches(1.25), SLIDE_W, Pt(4), fill=COPPER)
    tb, tf = textbox(slide, Inches(0.55), Inches(0.18), Inches(11.5), Inches(0.75), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p, title, 30, WHITE, bold=True, font='Georgia')
    if subtitle:
        tb2, tf2 = textbox(slide, Inches(0.55), Inches(0.92), Inches(10), Inches(0.3))
        set_run(tf2.paragraphs[0], subtitle, 13, RGBColor(0xE6, 0xD8, 0xC8), italic=True)
    # slide number chip
    tb3, tf3 = textbox(slide, SLIDE_W - Inches(1.0), Inches(0.4), Inches(0.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    set_run(p3, str(num), 16, COPPER, bold=True)


def footer(slide, num):
    tb, tf = textbox(slide, Inches(0.55), SLIDE_H - Inches(0.4), Inches(6), Inches(0.3))
    set_run(tf.paragraphs[0], "LENDORA  |  Rental Marketplace Defense", 10, INK_LT)
    tb2, tf2 = textbox(slide, SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.4), Inches(0.7), Inches(0.3))
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_run(p, f"{num} / 39", 10, INK_LT)


def bullets_block(slide, l, t, w, h, bullets, size=17, gap=10, color=INK, marker_color=COPPER):
    tb, tf = textbox(slide, l, t, w, h)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        set_run(p, "●  ", size - 2, marker_color, bold=True)
        set_run(p, b, size, color)
    return tb


def add_notes(slide, script, transition):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = script
    if transition:
        p = notes_tf.add_paragraph()
        p.text = f"\n➡ TRANSITION: {transition}"


# ---------------------------------------------------------------- title ----
def build_title(s):
    slide = add_slide()
    fill_bg(slide, FOREST)
    rect(slide, 0, Inches(3.55), SLIDE_W, Pt(3), fill=COPPER)
    tb, tf = textbox(slide, Inches(1), Inches(2.3), Inches(11.33), Inches(1.1), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, s['title'], 60, WHITE, bold=True, font='Georgia')
    tb2, tf2 = textbox(slide, Inches(1), Inches(3.75), Inches(11.33), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, s['subtitle'], 20, RGBColor(0xE6, 0xD8, 0xC8), italic=True)
    tb3, tf3 = textbox(slide, Inches(1), Inches(4.7), Inches(11.33), Inches(2.0))
    for i, m in enumerate(s['meta']):
        p3 = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_after = Pt(6)
        set_run(p3, m, 14, CREAM)
    add_notes(slide, s['script'], s['transition'])


# -------------------------------------------------------------- content ----
def build_content(s):
    slide = add_slide()
    fill_bg(slide, CREAM)
    header_bar(slide, s['num'], s['title'])
    bullets_block(slide, Inches(0.7), Inches(1.75), Inches(11.9), Inches(5), s['bullets'], size=19, gap=16)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------------ screenshot ---
def build_screenshot(s):
    slide = add_slide()
    fill_bg(slide, CREAM)
    header_bar(slide, s['num'], s['title'])

    left_w = Inches(5.0)
    left_l = Inches(0.55)
    top = Inches(1.6)

    def label_para(tf, label, body, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.space_after = Pt(2)
        set_run(p, label, 13, COPPER, bold=True)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(14)
        set_run(p2, body, 14.5, INK)

    tb, tf = textbox(slide, left_l, top, left_w, Inches(5.4))
    label_para(tf, "WHAT IT DOES", s['what'], first=True)
    label_para(tf, "WHO USES IT", s['who'])
    label_para(tf, "KEY ACTIONS", s['actions'])
    label_para(tf, "LEADS TO →", s['connects'])

    # screenshot placeholder(s) on the right
    right_l = Inches(5.85)
    right_w = SLIDE_W - right_l - Inches(0.55)
    n = len(s['shots'])
    gap = Inches(0.25)
    total_h = Inches(5.4)
    each_h = Emu(int((total_h - gap * (n - 1)) / n)) if n > 1 else total_h
    y = top
    for shot_label in s['shots']:
        ph = dashed_rect(slide, right_l, y, right_w, each_h, fill=PH_BG, line=COPPER)
        tf_ph = ph.text_frame
        tf_ph.word_wrap = True
        tf_ph.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_ph.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, "\U0001F4F7  INSERT SCREENSHOT", 15, COPPER, bold=True)
        p2 = tf_ph.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2, shot_label, 13, INK_LT, italic=True)
        y = Emu(int(y) + int(each_h) + int(gap))

    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ----------------------------------------------------- architecture diagram
def build_arch(s):
    slide = add_slide()
    fill_bg(slide, CREAM)
    header_bar(slide, s['num'], s['title'])

    box_h = Inches(0.95)
    y = Inches(1.85)
    labels = [
        ("Browser\n(Client)", FOREST),
        ("Next.js Frontend\nReact + TanStack Query", COPPER),
        ("REST API (Express)\nRoutes → Use Cases → Repositories", FOREST),
        ("MySQL Database\n(via Prisma ORM)", COPPER),
    ]
    n = len(labels)
    gap = Inches(0.35)
    box_w = Emu(int((SLIDE_W - Inches(1.1) - gap * (n - 1)) / n))
    x = Inches(0.55)
    centers = []
    for text, color in labels:
        sp = rect(slide, x, y, box_w, box_h, fill=color)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        lines = text.split("\n")
        set_run(p, lines[0], 14, WHITE, bold=True)
        for ln in lines[1:]:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            set_run(p2, ln, 10.5, RGBColor(0xEE, 0xE6, 0xDC))
        centers.append((Emu(int(x) + int(box_w) // 2), int(x), int(box_w)))
        x = Emu(int(x) + int(box_w) + int(gap))

    # arrows between boxes
    arrow_y = Emu(int(y) + int(box_h) // 2)
    bx = Inches(0.55)
    for i in range(n - 1):
        ax1 = Emu(int(bx) + int(box_w))
        ax2 = Emu(int(ax1) + int(gap))
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax1, arrow_y, ax2, arrow_y)
        conn.line.color.rgb = INK_LT
        conn.line.width = Pt(2.25)
        conn.line._get_or_add_ln().append(
            conn.line._get_or_add_ln().makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        bx = Emu(int(bx) + int(box_w) + int(gap))

    # side attachments: Redis + Gateway under API box, Nginx above all
    api_x = Emu(int(Inches(0.55)) + int(box_w) * 2 + int(gap) * 2)
    side_y = Emu(int(y) + int(box_h) + Inches(0.55))
    for i, (txt, col) in enumerate([("Redis Cache", COPPER), ("bKash / Nagad\nGateway", FOREST)]):
        sx = Emu(int(api_x) + i * (int(box_w) // 2 + 20))
        sp = rect(slide, sx, side_y, Emu(int(box_w) // 2 - 10), Inches(0.75), fill=col)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        lines = txt.split("\n")
        set_run(p, lines[0], 11, WHITE, bold=True)
        for ln in lines[1:]:
            p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
            set_run(p2, ln, 11, WHITE, bold=True)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                           Emu(int(sx) + int(box_w) // 4 - 5), side_y,
                                           Emu(int(api_x) + int(box_w) // 2), Emu(int(y) + int(box_h)))
        conn.line.color.rgb = INK_LT
        conn.line.width = Pt(1.5)

    bullets_block(slide, Inches(0.55), Inches(4.55), Inches(12.2), Inches(2.4), s['bullets'], size=15.5, gap=8)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# --------------------------------------------------------- workflow diagram
def build_workflow(s):
    slide = add_slide()
    fill_bg(slide, CREAM)
    header_bar(slide, s['num'], s['title'])

    steps = s['steps']
    n = len(steps)
    cols = 5
    rows = 2
    gap_x = Inches(0.22)
    gap_y = Inches(0.7)
    box_w = Emu(int((SLIDE_W - Inches(1.1) - gap_x * (cols - 1)) / cols))
    box_h = Inches(1.05)
    start_x = Inches(0.55)
    start_y = Inches(2.0)
    colors = [FOREST, COPPER]

    positions = []
    for i, step in enumerate(steps):
        row = i // cols
        col = i % cols
        x = Emu(int(start_x) + col * (int(box_w) + int(gap_x)))
        y = Emu(int(start_y) + row * (int(box_h) + int(gap_y)))
        positions.append((x, y))

    for i, step in enumerate(steps):
        x, y = positions[i]
        col_color = colors[i % 2]
        sp = rect(slide, x, y, box_w, box_h, fill=col_color)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        lines = step.split("\n")
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            set_run(p, ln, 12.5, WHITE, bold=True)
        # step number badge
        bx, by = Emu(int(x) - 14), Emu(int(y) - 14)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, Inches(0.34), Inches(0.34))
        badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
        badge.line.color.rgb = INK; badge.line.width = Pt(1)
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        set_run(bp, str(i + 1), 11, INK, bold=True)
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # arrows: row1 left-to-right, drop to row2, row2 left-to-right
    def arrow(x1, y1, x2, y2):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = INK_LT
        conn.line.width = Pt(2)

    for col in range(cols - 1):
        x1, y1 = positions[col]
        x2, y2 = positions[col + 1]
        midy = Emu(int(y1) + int(box_h) // 2)
        arrow(Emu(int(x1) + int(box_w)), midy, x2, midy)
    # drop from end of row1 to start of row2
    x1, y1 = positions[cols - 1]
    x2, y2 = positions[cols]
    arrow(Emu(int(x1) + int(box_w) // 2), Emu(int(y1) + int(box_h)), Emu(int(x2) + int(box_w) // 2), y2)
    for col in range(cols, n - 1):
        x1, y1 = positions[col]
        x2, y2 = positions[col + 1]
        midy = Emu(int(y1) + int(box_h) // 2)
        arrow(Emu(int(x1) + int(box_w)), midy, x2, midy)

    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------------- db diagram --
def build_db(s):
    slide = add_slide()
    fill_bg(slide, CREAM)
    header_bar(slide, s['num'], s['title'])

    tables = ["User", "VendorProfile", "Product", "Rental", "Payment",
              "VendorPayout", "ReturnRecord", "LateFeeTransaction", "Review", "Dispute"]
    cols = 5
    box_w = Inches(2.2)
    box_h = Inches(0.65)
    gap_x = Inches(0.25)
    gap_y = Inches(0.45)
    start_x = Inches(0.65)
    start_y = Inches(1.75)
    hub_idx = 3  # Rental
    positions = []
    for i, t in enumerate(tables):
        row = i // cols
        col = i % cols
        x = Emu(int(start_x) + col * (int(box_w) + int(gap_x)))
        y = Emu(int(start_y) + row * (int(box_h) + int(gap_y)))
        positions.append((x, y))

    for i, t in enumerate(tables):
        x, y = positions[i]
        color = COPPER if i == hub_idx else FOREST
        sp = rect(slide, x, y, box_w, box_h, fill=color)
        tf = sp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, t, 13, WHITE, bold=True)

    hx, hy = positions[hub_idx]
    hub_center = (Emu(int(hx) + int(box_w) // 2), Emu(int(hy) + int(box_h) // 2))
    for i in [0, 1, 2, 4, 5, 6, 7, 8, 9]:
        x, y = positions[i]
        cx, cy = Emu(int(x) + int(box_w) // 2), Emu(int(y) + int(box_h) // 2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, hub_center[0], hub_center[1], cx, cy)
        conn.line.color.rgb = RGBColor(0xC9, 0xBF, 0xAF)
        conn.line.width = Pt(1.25)
    # redraw hub on top so lines don't overlap text
    sp = rect(slide, hx, hy, box_w, box_h, fill=COPPER)
    tf = sp.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Rental (hub)", 13, WHITE, bold=True)

    bullets_block(slide, Inches(0.65), Inches(4.55), Inches(12.0), Inches(2.4), s['bullets'], size=14.5, gap=7)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------------------ thanks -
def build_thanks(s):
    slide = add_slide()
    fill_bg(slide, FOREST)
    rect(slide, 0, Inches(3.3), SLIDE_W, Pt(3), fill=COPPER)
    tb, tf = textbox(slide, Inches(1), Inches(2.2), Inches(11.33), Inches(1.0), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, s['title'], 54, WHITE, bold=True, font='Georgia')
    tb2, tf2 = textbox(slide, Inches(1), Inches(3.5), Inches(11.33), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    set_run(p2, s['subtitle'], 20, RGBColor(0xE6, 0xD8, 0xC8), italic=True)
    tb3, tf3 = textbox(slide, Inches(1), Inches(4.5), Inches(11.33), Inches(1.6))
    for i, m in enumerate(s['meta']):
        p3 = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_after = Pt(6)
        set_run(p3, m, 14, CREAM)
    add_notes(slide, s['script'], s['transition'])


BUILDERS = {
    'title': build_title,
    'content': build_content,
    'screenshot': build_screenshot,
    'diagram_arch': build_arch,
    'diagram_workflow': build_workflow,
    'diagram_db': build_db,
    'thanks': build_thanks,
}

for s in SLIDES:
    BUILDERS[s['kind']](s)

out_path = "Lendora_Defense_Presentation.pptx"
prs.save(out_path)
print(f"Saved {out_path} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
