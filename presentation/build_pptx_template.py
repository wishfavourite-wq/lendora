# -*- coding: utf-8 -*-
"""Rebuilds the Lendora defense deck on top of elegant-workplan.pptx, preserving
its theme/fonts/decorative floral layouts, and adds a matching icon set.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from slides_data import SLIDES

DK1    = RGBColor(0x1C, 0x34, 0x3C)   # near-black teal — primary text
CREAM  = RGBColor(0xF0, 0xEB, 0xE4)   # background
DK2    = RGBColor(0x4C, 0x6A, 0x78)   # blue-grey accent
SAGE   = RGBColor(0xA7, 0xB3, 0xB2)   # sage accent
TERRA  = RGBColor(0xE0, 0xB4, 0xA4)   # terracotta accent
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PH_BG  = RGBColor(0xE3, 0xDC, 0xD2)

TITLE_FONT = 'Poppins'
BODY_FONT  = 'Arimo'

ICON_DIR = "icons"

TOPIC_ICON = {
    2: 'lightbulb', 3: 'alert', 4: 'target', 5: 'x_circle', 6: 'check_circle',
    10: 'house', 11: 'store', 12: 'user_plus', 13: 'shield_check', 14: 'grid',
    15: 'box', 16: 'grid', 17: 'search', 18: 'calendar', 19: 'check_circle',
    20: 'wallet', 21: 'truck', 22: 'alert', 23: 'refresh', 24: 'camera',
    25: 'wallet', 26: 'percent', 27: 'star', 28: 'chart', 30: 'lock',
    32: 'lightning', 33: 'wrench', 34: 'minus_circle', 35: 'rocket',
    36: 'clipboard_check', 37: 'trophy', 38: 'flag',
}
MINI_ICONS = ('document', 'user', 'lightning', 'arrow_right')

prs = Presentation("elegant-workplan.pptx")
SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height

LAYOUTS = prs.slide_masters[0].slide_layouts
L_TITLE = LAYOUTS[0]
L_TITLE_ONLY = LAYOUTS[4]


def delete_all_slides():
    xml_slides = prs.slides._sldIdLst
    for sldId in list(xml_slides):
        rId = sldId.get(qn('r:id'))
        prs.part.drop_rel(rId)
        xml_slides.remove(sldId)


def icon_path(key):
    return f"{ICON_DIR}/{key}.png"


def textbox(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
        setattr(tf, attr, 0)
    return tb, tf


def set_run(p, text, size, color, bold=False, italic=False, font=BODY_FONT):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None, dashed=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sp.adjustments[0] = 0.06
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w or Pt(1)
        if dashed:
            ln = sp.line._get_or_add_ln()
            ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    return sp


def add_icon(slide, key, l, t, size):
    try:
        slide.shapes.add_picture(icon_path(key), l, t, width=size, height=size)
    except Exception:
        pass


def footer(slide, num):
    # Centered, bottom-middle — both bottom corners carry decorative floral artwork
    tb, tf = textbox(slide, Inches(2.3), SLIDE_H - Inches(0.32), Inches(5.4), Inches(0.25))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, f"LENDORA  •  Rental Marketplace Defense  •  Slide {num}/39", 8.5, DK2)


def add_notes(slide, script, transition):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = script
    if transition:
        p = notes_tf.add_paragraph()
        p.text = f"\n➡ TRANSITION: {transition}"


def title_only_slide():
    return prs.slides.add_slide(L_TITLE_ONLY)


def set_title(slide, text, size=None):
    ph = slide.placeholders[0]
    ph.text_frame.text = text
    if size:
        for p in ph.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(size)


def bullets_block(slide, l, t, w, h, bullets, size=13, gap=10):
    tb, tf = textbox(slide, l, t, w, h)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        set_run(p, "●  ", size - 2, TERRA, bold=True)
        set_run(p, b, size, DK1)
    return tb


# ------------------------------------------------------------------- title -
def build_title(s):
    slide = prs.slides.add_slide(L_TITLE)
    ph_title = slide.placeholders[0]
    ph_title.text_frame.text = s['title']
    for p in ph_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(46)
            r.font.color.rgb = DK1
    ph_sub = slide.placeholders[1]
    ph_sub.text_frame.text = s['subtitle']
    for p in ph_sub.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(15)
            r.font.name = BODY_FONT
            r.font.color.rgb = DK2
            r.font.italic = True

    tb, tf = textbox(slide, Inches(1.6), Inches(4.25), Inches(6.8), Inches(1.1), anchor=MSO_ANCHOR.TOP)
    for i, m in enumerate(s['meta']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        set_run(p, m, 10.5, DK1)
    add_notes(slide, s['script'], s['transition'])


# ----------------------------------------------------------------- thanks --
def build_thanks(s):
    slide = prs.slides.add_slide(L_TITLE)
    ph_title = slide.placeholders[0]
    ph_title.text_frame.text = s['title']
    for p in ph_title.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(48)
            r.font.color.rgb = DK1
    ph_sub = slide.placeholders[1]
    ph_sub.text_frame.text = s['subtitle']
    for p in ph_sub.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.name = BODY_FONT
            r.font.color.rgb = DK2
            r.font.italic = True

    add_icon(slide, 'heart', Inches(4.5), Inches(4.3), Inches(1.0))
    tb, tf = textbox(slide, Inches(1.6), Inches(4.25), Inches(6.8), Inches(0.9), anchor=MSO_ANCHOR.TOP)
    for i, m in enumerate(s['meta']):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        set_run(p, m, 11, DK1)
    add_notes(slide, s['script'], s['transition'])


# --------------------------------------------------------------- content ---
def build_content(s):
    slide = title_only_slide()
    set_title(slide, s['title'])
    icon_key = TOPIC_ICON.get(s['num'])
    if icon_key:
        add_icon(slide, icon_key, SLIDE_W - Inches(1.55), Inches(0.6), Inches(0.5))
    bullets_block(slide, Inches(1.3), Inches(1.55), Inches(7.75), Inches(3.7), s['bullets'], size=13.5, gap=11)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ----------------------------------------------------------- tech stack ----
TECH_ICONS = ['monitor', 'server', 'database', 'card', 'layers']

def build_tech(s):
    slide = title_only_slide()
    set_title(slide, s['title'])
    top = Inches(1.55)
    row_h = Inches(0.78)
    for i, b in enumerate(s['bullets']):
        y = Emu(int(top) + i * int(row_h))
        ic = TECH_ICONS[i] if i < len(TECH_ICONS) else 'check_circle'
        add_icon(slide, ic, Inches(1.3), y, Inches(0.42))
        tb, tf = textbox(slide, Inches(1.85), y, Inches(7.25), row_h, anchor=MSO_ANCHOR.MIDDLE)
        set_run(tf.paragraphs[0], b, 12.5, DK1)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ----------------------------------------------------------- user roles ---
def build_roles(s):
    slide = title_only_slide()
    set_title(slide, s['title'])
    roles = [('user', s['bullets'][0]), ('store', s['bullets'][1]), ('shield_check', s['bullets'][2])]
    col_w = Inches(2.85)
    gap = Inches(0.2)
    x = Inches(0.55)
    top = Inches(1.7)
    for icon_key, text in roles:
        rect(slide, x, top, col_w, Inches(2.55), fill=WHITE, line=SAGE, line_w=Pt(1))
        add_icon(slide, icon_key, Emu(int(x) + int(col_w) // 2 - int(Inches(0.32))), Emu(int(top) + Inches(0.22)), Inches(0.64))
        tb, tf = textbox(slide, Emu(int(x) + Inches(0.18)), Emu(int(top) + Inches(1.05)), Emu(int(col_w) - Inches(0.36)), Inches(1.4))
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, text, 11, DK1)
        x = Emu(int(x) + int(col_w) + int(gap))
    tb2, tf2 = textbox(slide, Inches(1.3), Inches(4.45), Inches(7.75), Inches(0.6))
    p_note = tf2.paragraphs[0]
    p_note.alignment = PP_ALIGN.CENTER
    set_run(p_note, s['bullets'][3], 11, DK2, italic=True)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------------- screenshot --
def build_screenshot(s):
    slide = title_only_slide()
    set_title(slide, s['title'])
    icon_key = TOPIC_ICON.get(s['num'])
    if icon_key:
        add_icon(slide, icon_key, SLIDE_W - Inches(1.55), Inches(0.6), Inches(0.5))

    left_l = Inches(1.3)
    left_w = Inches(3.05)
    top = Inches(1.55)
    rows = [("WHAT IT DOES", s['what']), ("WHO USES IT", s['who']),
            ("KEY ACTIONS", s['actions']), ("LEADS TO", s['connects'])]
    y = top
    for i, (label, body) in enumerate(rows):
        add_icon(slide, MINI_ICONS[i], left_l, y, Inches(0.26))
        tb, tf = textbox(slide, Emu(int(left_l) + Inches(0.34)), y, Emu(int(left_w) - Inches(0.34)), Inches(0.9))
        p = tf.paragraphs[0]
        set_run(p, label, 9.5, TERRA, bold=True)
        p2 = tf.add_paragraph()
        p2.space_before = Pt(1)
        set_run(p2, body, 10.8, DK1)
        y = Emu(int(y) + Inches(0.92))

    right_l = Inches(4.65)
    right_w = Inches(9.15) - right_l  # stop short of the bottom-right decorative dot
    shots = s['shots']
    n = len(shots)
    gap = Inches(0.18)
    total_h = Inches(3.55)
    each_h = Emu(int((total_h - gap * (n - 1)) / n)) if n > 1 else total_h
    y = top
    for shot_label in shots:
        ph = rect(slide, right_l, y, right_w, each_h, fill=PH_BG, line=DK2, line_w=Pt(1.5), dashed=True)
        tf_ph = ph.text_frame
        tf_ph.word_wrap = True
        tf_ph.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf_ph.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, "\U0001F4F7  INSERT SCREENSHOT", 12, DK2, bold=True)
        p2 = tf_ph.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2, shot_label, 10, DK2, italic=True)
        y = Emu(int(y) + int(each_h) + int(gap))

    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------- architecture -----
def build_arch(s):
    slide = title_only_slide()
    set_title(slide, s['title'])

    box_h = Inches(0.72)
    y = Inches(1.55)
    labels = [("Browser", 'monitor', SAGE), ("Next.js Frontend", 'layers', TERRA),
              ("REST API (Express)", 'server', DK2), ("MySQL (Prisma)", 'database', SAGE)]
    n = len(labels)
    gap = Inches(0.22)
    box_w = Emu(int((SLIDE_W - Inches(0.9) - gap * (n - 1)) / n))
    x = Inches(0.45)
    box_x = []
    for text, icon_key, color in labels:
        sp = rect(slide, x, y, box_w, box_h, fill=color)
        add_icon(slide, icon_key, Emu(int(x) + int(box_w) // 2 - Inches(0.13)), Emu(int(y) - Inches(0.42)), Inches(0.36))
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, text, 10, WHITE if color != SAGE else DK1, bold=True)
        box_x.append(int(x))
        x = Emu(int(x) + int(box_w) + int(gap))

    arrow_y = Emu(int(y) + int(box_h) // 2)
    for i in range(n - 1):
        ax1 = Emu(box_x[i] + int(box_w))
        ax2 = Emu(box_x[i + 1])
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, ax1, arrow_y, ax2, arrow_y)
        conn.line.color.rgb = DK2
        conn.line.width = Pt(2)
        ln = conn.line._get_or_add_ln()
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))

    side_y = Emu(int(y) + int(box_h) + Inches(0.45))
    api_x = box_x[2]
    for i, (txt, icon_key) in enumerate([("Redis Cache", 'database'), ("bKash / Nagad", 'wallet')]):
        sx = Emu(api_x + i * (int(box_w) // 2 + 8))
        sw = Emu(int(box_w) // 2 - 8)
        sp = rect(slide, sx, side_y, sw, Inches(0.6), fill=DK1)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, txt, 9, WHITE, bold=True)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                           Emu(int(sx) + int(sw) // 2), side_y,
                                           Emu(api_x + int(box_w) // 2), Emu(int(y) + int(box_h)))
        conn.line.color.rgb = DK2
        conn.line.width = Pt(1.25)

    bullets_block(slide, Inches(1.3), Inches(3.15), Inches(7.8), Inches(2.2), s['bullets'], size=11, gap=6)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# ------------------------------------------------------------- workflow ----
def build_workflow(s):
    slide = title_only_slide()
    set_title(slide, s['title'])

    steps = s['steps']
    cols = 5
    gap_x = Inches(0.16)
    gap_y = Inches(0.55)
    box_w = Emu(int((SLIDE_W - Inches(0.9) - gap_x * (cols - 1)) / cols))
    box_h = Inches(0.85)
    start_x = Inches(0.45)
    start_y = Inches(1.65)
    colors = [DK2, TERRA]

    positions = []
    for i in range(len(steps)):
        row, col = divmod(i, cols)
        x = Emu(int(start_x) + col * (int(box_w) + int(gap_x)))
        y = Emu(int(start_y) + row * (int(box_h) + int(gap_y)))
        positions.append((x, y))

    for i, step in enumerate(steps):
        x, y = positions[i]
        color = colors[i % 2]
        sp = rect(slide, x, y, box_w, box_h, fill=color)
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for j, ln in enumerate(step.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            set_run(p, ln, 9.5, WHITE, bold=True)
        bx, by = Emu(int(x) - Inches(0.1)), Emu(int(y) - Inches(0.1))
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, bx, by, Inches(0.26), Inches(0.26))
        badge.fill.solid(); badge.fill.fore_color.rgb = WHITE
        badge.line.color.rgb = DK1; badge.line.width = Pt(1)
        badge.shadow.inherit = False
        bp = badge.text_frame.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        set_run(bp, str(i + 1), 9, DK1, bold=True)
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    def arrow(x1, y1, x2, y2):
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        conn.line.color.rgb = DK2
        conn.line.width = Pt(1.75)

    for col in range(cols - 1):
        x1, y1 = positions[col]; x2, _ = positions[col + 1]
        midy = Emu(int(y1) + int(box_h) // 2)
        arrow(Emu(int(x1) + int(box_w)), midy, x2, midy)
    x1, y1 = positions[cols - 1]; x2, y2 = positions[cols]
    arrow(Emu(int(x1) + int(box_w) // 2), Emu(int(y1) + int(box_h)), Emu(int(x2) + int(box_w) // 2), y2)
    for col in range(cols, len(steps) - 1):
        x1, y1 = positions[col]; x2, _ = positions[col + 1]
        midy = Emu(int(y1) + int(box_h) // 2)
        arrow(Emu(int(x1) + int(box_w)), midy, x2, midy)

    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


# -------------------------------------------------------------- db design --
def build_db(s):
    slide = title_only_slide()
    set_title(slide, s['title'])

    tables = ["User", "VendorProfile", "Product", "Rental", "Payment",
              "VendorPayout", "ReturnRecord", "LateFeeTransaction", "Review", "Dispute"]
    cols = 5
    box_w = Inches(1.7)
    box_h = Inches(0.52)
    gap_x = Inches(0.2)
    gap_y = Inches(0.35)
    start_x = Inches(0.55)
    start_y = Inches(1.5)
    hub_idx = 3
    positions = []
    for i in range(len(tables)):
        row, col = divmod(i, cols)
        x = Emu(int(start_x) + col * (int(box_w) + int(gap_x)))
        y = Emu(int(start_y) + row * (int(box_h) + int(gap_y)))
        positions.append((x, y))

    hx, hy = positions[hub_idx]
    hub_center = (Emu(int(hx) + int(box_w) // 2), Emu(int(hy) + int(box_h) // 2))
    for i in range(len(tables)):
        if i == hub_idx:
            continue
        x, y = positions[i]
        cx, cy = Emu(int(x) + int(box_w) // 2), Emu(int(y) + int(box_h) // 2)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, hub_center[0], hub_center[1], cx, cy)
        conn.line.color.rgb = SAGE
        conn.line.width = Pt(1)

    for i, t in enumerate(tables):
        x, y = positions[i]
        color = TERRA if i == hub_idx else DK2
        sp = rect(slide, x, y, box_w, box_h, fill=color)
        tf = sp.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, "Rental (hub)" if i == hub_idx else t, 10, DK1 if color == TERRA else WHITE, bold=True)

    bullets_block(slide, Inches(1.3), Inches(3.95), Inches(7.8), Inches(1.5), s['bullets'], size=10.5, gap=5)
    footer(slide, s['num'])
    add_notes(slide, s['script'], s['transition'])


BUILDERS = {
    'title': build_title,
    'content': build_content,
    'screenshot': build_screenshot,
    'diagram_arch': build_arch,
    'diagram_workflow': build_workflow,
    'diagram_db': build_db,
    'thanks': build_thanks,
}

delete_all_slides()

for s in SLIDES:
    if s['num'] == 8:
        build_tech(s)
    elif s['num'] == 9:
        build_roles(s)
    else:
        BUILDERS[s['kind']](s)

out_path = "Lendora_Defense_Presentation.pptx"
prs.save(out_path)
print(f"Saved {out_path}")
